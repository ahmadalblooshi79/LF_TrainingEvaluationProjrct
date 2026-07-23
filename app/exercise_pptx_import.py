"""استيراد أقسام التمرين من عرض PowerPoint (.pptx).

يشمل: معلومات التمرين (القصد، المشاركون، نوع/مستوى)، الفكرة العامة/الخاصة،
البرنامج، والخريطة.
"""

from __future__ import annotations

import re
import unicodedata
import zipfile
from io import BytesIO

# مفاتيح التبويبات في صفحة التمرين الحالي
TAB_KEYS = ("general", "specific", "program", "map")

# حقول لوحة «معلومات التمرين»
INFO_KEYS = ("purpose", "participants", "type_level")

TAB_LABELS: dict[str, tuple[str, ...]] = {
    "general": ("الفكرة العامة", "فكرة عامة"),
    "specific": ("الفكرة الخاصة", "فكرة خاصة"),
    "program": ("البرنامج", "برنامج التمرين", "برنامج"),
    "map": ("الخريطة", "خريطة التمرين", "خريطة"),
}

INFO_LABELS: dict[str, tuple[str, ...]] = {
    "purpose": ("القصد من التمرين", "القصد"),
    "participants": (
        "المشاركون في التمرين",
        "المشاركين في التمرين",
        "المشاركون",
        "المشاركين",
    ),
    "type_level": ("نوع ومستوى التمرين", "نوع ومستوى"),
}

# شريحة جامعة تُقسَّم داخلياً إلى حقول المعلومات
INFO_SLIDE_LABELS = ("معلومات التمرين", "معلومات")

FIELD_BY_TAB: dict[str, str] = {
    "general": "general_idea_text",
    "specific": "specific_idea_text",
    "program": "program_text",
    "map": "map_text",
}

FIELD_BY_INFO: dict[str, str] = {
    "purpose": "exercise_purpose",
    "participants": "exercise_participants",
    "type_level": "exercise_type_level_text",
}

SECTION_LABELS: dict[str, tuple[str, ...]] = {**TAB_LABELS, **INFO_LABELS}
FIELD_BY_SECTION: dict[str, str] = {**FIELD_BY_TAB, **FIELD_BY_INFO}

_AR_TATWEEL = "\u0640"
_NORM_RE = re.compile(r"\s+")


def is_pptx_bytes(data: bytes) -> bool:
    """ملف PowerPoint .pptx — أرشيف ZIP يحتوي مجلد ppt/."""
    if not data or len(data) < 64:
        return False
    if data[:2] != b"PK":
        return False
    try:
        with zipfile.ZipFile(BytesIO(data)) as zf:
            return any(n.startswith("ppt/") for n in zf.namelist())
    except zipfile.BadZipFile:
        return False


def _normalize_label(text: str) -> str:
    s = unicodedata.normalize("NFKC", (text or "").strip())
    s = s.replace(_AR_TATWEEL, "")
    s = _NORM_RE.sub(" ", s)
    # توحيد أشكال الألف
    for ch in ("أ", "إ", "آ"):
        s = s.replace(ch, "ا")
    s = s.replace("ى", "ي")
    return s.strip().lower()


def _match_section_key(text: str) -> str | None:
    norm = _normalize_label(text)
    if not norm:
        return None
    best_key: str | None = None
    best_len = 0
    for key, labels in SECTION_LABELS.items():
        for label in labels:
            ln = _normalize_label(label)
            if not ln:
                continue
            if norm == ln or ln in norm or norm in ln:
                if len(ln) > best_len:
                    best_key = key
                    best_len = len(ln)
    return best_key


def _match_tab_key(text: str) -> str | None:
    key = _match_section_key(text)
    return key if key in TAB_LABELS else None


def _is_info_slide_title(text: str) -> bool:
    norm = _normalize_label(text)
    if not norm:
        return False
    for label in INFO_SLIDE_LABELS:
        ln = _normalize_label(label)
        if norm == ln or ln in norm:
            return True
    return False


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts: list[str] = []
    for para in shape.text_frame.paragraphs:
        line = "".join((run.text or "") for run in para.runs).strip()
        if not line:
            line = (para.text or "").strip()
        if line:
            parts.append(line)
    return "\n".join(parts).strip()


def _slide_title_text(slide) -> str:
    try:
        if slide.shapes.title is not None:
            return _shape_text(slide.shapes.title)
    except Exception:
        pass
    return ""


def _slide_program_key(slide) -> bool:
    title = _slide_title_text(slide)
    if title and _match_tab_key(title) == "program":
        return True
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            continue
        txt = _shape_text(shape)
        if txt and _match_tab_key(txt.split("\n", 1)[0]) == "program":
            return True
    return False


def _slide_body_texts(slide, *, skip_title: str) -> list[str]:
    skip_norm = _normalize_label(skip_title)
    out: list[str] = []
    for shape in slide.shapes:
        txt = _shape_text(shape)
        if not txt:
            continue
        if skip_norm and _normalize_label(txt) == skip_norm:
            continue
        if shape is slide.shapes.title:
            continue
        out.append(txt)
    return out


def _merge_field(fields: dict[str, str], key: str, chunk: str) -> None:
    chunk = (chunk or "").strip()
    if not chunk:
        return
    field = FIELD_BY_SECTION[key]
    if fields[field]:
        if chunk not in fields[field]:
            fields[field] = f"{fields[field].rstrip()}\n\n{chunk}"
    else:
        fields[field] = chunk


def _parse_sections_in_text(text: str, fields: dict[str, str]) -> None:
    current: str | None = None
    buf: list[str] = []
    for raw_line in (text or "").splitlines():
        line = raw_line.strip()
        if not line:
            if buf:
                buf.append("")
            continue
        key = _match_section_key(line)
        if key:
            if current and buf:
                _merge_field(fields, current, "\n".join(buf).strip())
            current = key
            buf = []
            continue
        if current:
            buf.append(raw_line.rstrip())
    if current and buf:
        _merge_field(fields, current, "\n".join(buf).strip())


def parse_exercise_pptx_bytes(data: bytes) -> dict:
    """استخراج نصوص التبويبات وحقول المعلومات وجدول البرنامج من ملف .pptx."""
    empty_fields = {FIELD_BY_TAB[k]: "" for k in TAB_KEYS}
    empty_fields.update({FIELD_BY_INFO[k]: "" for k in INFO_KEYS})
    empty_fields["program_table_json"] = ""
    if not is_pptx_bytes(data):
        return {"ok": False, "error": "bad_pptx", "fields": empty_fields, "warnings": []}
    try:
        from pptx import Presentation

        from app.exercise_program_table import (
            dumps_program_table,
            extract_program_table_from_slide,
        )
    except ImportError:
        return {
            "ok": False,
            "error": "missing_dependency",
            "fields": empty_fields,
            "warnings": [],
        }
    try:
        prs = Presentation(BytesIO(data))
    except Exception:
        return {"ok": False, "error": "invalid_pptx", "fields": empty_fields, "warnings": []}

    fields = dict(empty_fields)
    warnings: list[str] = []

    for slide in prs.slides:
        title = _slide_title_text(slide)
        if _slide_program_key(slide):
            table_data = extract_program_table_from_slide(slide)
            if table_data:
                fields["program_table_json"] = dumps_program_table(table_data)
                continue

        title_key = _match_section_key(title) if title else None
        skip_title = title if (title_key or _is_info_slide_title(title)) else ""
        bodies = _slide_body_texts(slide, skip_title=skip_title)
        body_joined = "\n\n".join(bodies).strip()

        if title and _is_info_slide_title(title):
            full_slide = "\n\n".join([t for t in [body_joined] if t]).strip()
            if full_slide:
                _parse_sections_in_text(full_slide, fields)
            continue

        if title_key:
            _merge_field(fields, title_key, body_joined or title)
            continue

        full_slide = "\n\n".join([t for t in [title, body_joined] if t]).strip()
        if full_slide:
            _parse_sections_in_text(full_slide, fields)
            if not title and bodies:
                first_key = _match_section_key(bodies[0].split("\n", 1)[0])
                if first_key and len(bodies) == 1:
                    rest = "\n".join(bodies[0].splitlines()[1:]).strip()
                    if rest:
                        _merge_field(fields, first_key, rest)

    filled_tabs = [
        k
        for k in TAB_KEYS
        if (fields[FIELD_BY_TAB[k]] or "").strip()
        or (k == "program" and (fields.get("program_table_json") or "").strip())
    ]
    filled_info = [
        k for k in INFO_KEYS if (fields[FIELD_BY_INFO[k]] or "").strip()
    ]
    if not filled_tabs and not filled_info:
        return {
            "ok": False,
            "error": "no_content",
            "fields": fields,
            "warnings": warnings,
        }
    for key in TAB_KEYS:
        if key == "program" and (fields.get("program_table_json") or "").strip():
            continue
        if not (fields[FIELD_BY_TAB[key]] or "").strip():
            label = TAB_LABELS[key][0]
            warnings.append(f"لم يُعثر على محتوى لتبويب «{label}».")
    program_table_html = ""
    if fields.get("program_table_json"):
        from app.exercise_program_table import render_program_table_html

        program_table_html = render_program_table_html(fields["program_table_json"])
    return {
        "ok": True,
        "error": "",
        "fields": fields,
        "program_table_html": program_table_html,
        "warnings": warnings,
    }
