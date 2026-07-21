"""استخراج وعرض جدول برنامج التمرين من شريحة PowerPoint."""

from __future__ import annotations

import html
import json
import re
from typing import Any

_DATE_RE = re.compile(r"^\s*(\d{1,2})\s*/\s*(\d{1,2})\s*$")
_DATE_PREFIX_RE = re.compile(r"^(\d{1,2})\s*/\s*(\d{1,2})(?:\s+|$)")
_TATWEEL_RE = re.compile("\u0640+")
_AR_DAYS = (
    "الاثنين",
    "الإثنين",
    "الثلاثاء",
    "الأربعاء",
    "الخميس",
    "الجمعة",
    "السبت",
    "الأحد",
)


def _normalize_day(s: str) -> str:
    s = (s or "").strip()
    for ch in ("أ", "إ", "آ"):
        s = s.replace(ch, "ا")
    return s.replace("ى", "ي")


def _cell_fill_rgb(cell) -> tuple[int, int, int] | None:
    try:
        fill = cell.fill
        if fill.type is None:
            return None
        fc = fill.fore_color
        if fc.type is None:
            return None
        rgb = fc.rgb
        if rgb is None:
            return None
        return (int(rgb[0]), int(rgb[1]), int(rgb[2]))
    except Exception:
        return None


def _classify_week_band(rgb: tuple[int, int, int] | None) -> str:
    if rgb is None:
        return ""
    r, g, b = rgb
    if b >= 170 and b > r + 25 and b > g + 10:
        return "blue"
    if max(r, g, b) - min(r, g, b) < 35 and 150 <= (r + g + b) / 3 <= 230:
        return "grey"
    if r >= 210 and g >= 200 and b >= 170:
        return "beige"
    if r >= 180 and g >= 170 and b >= 140:
        return "beige"
    return ""


def _run_is_red(run) -> bool:
    try:
        color = run.font.color
        if color is None or color.type is None:
            return False
        rgb = color.rgb
        if rgb is None:
            return False
        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
        return r >= 140 and g <= 110 and b <= 110 and r > g + 30
    except Exception:
        return False


def _cell_text_parts(cell) -> list[dict[str, Any]]:
    parts: list[dict[str, Any]] = []
    try:
        tf = cell.text_frame
    except Exception:
        text = (getattr(cell, "text", None) or "").strip()
        if text:
            parts.append({"text": text, "red": False})
        return parts
    for para in tf.paragraphs:
        runs = list(para.runs)
        if runs:
            for run in runs:
                t = (run.text or "").strip()
                if not t:
                    continue
                parts.append({"text": t, "red": _run_is_red(run)})
        else:
            t = (para.text or "").strip()
            if t:
                parts.append({"text": t, "red": False})
    if not parts:
        t = (cell.text or "").strip()
        if t:
            parts.append({"text": t, "red": False})
    return parts


def _cell_colspan(cell) -> int:
    try:
        span = cell._tc.gridSpan  # noqa: SLF001
        if not span:
            return 1
        return int(span)
    except Exception:
        return 1


def _cell_rowspan(table, row_idx: int, col_idx: int) -> int:
    try:
        cell = table.cell(row_idx, col_idx)
        if cell.is_spanned:
            return 0
        v_merge = cell._tc.vMerge  # noqa: SLF001
        if not v_merge:
            return 1
        if str(v_merge) != "restart":
            return 0
        count = 1
        for ri in range(row_idx + 1, len(table.rows)):
            below = table.cell(ri, col_idx)
            if str(below._tc.vMerge) == "continue":  # noqa: SLF001
                count += 1
            else:
                break
        return count
    except Exception:
        return 1


def _format_date_label(month_or_day_a: str, month_or_day_b: str) -> str:
    """توحيد عرض التاريخ كـ m/d كما في جدول الأعمال."""
    a = int(month_or_day_a)
    b = int(month_or_day_b)
    return f"{a}/{b}"


def _split_date_and_parts(parts: list[dict[str, Any]]) -> tuple[str, list[dict[str, Any]]]:
    if not parts:
        return "", []
    first = _TATWEEL_RE.sub("", (parts[0].get("text") or "").strip())
    m = _DATE_RE.match(first)
    if m:
        return _format_date_label(m.group(1), m.group(2)), parts[1:]
    m2 = _DATE_PREFIX_RE.match(first)
    if m2:
        date_label = _format_date_label(m2.group(1), m2.group(2))
        rest = first[m2.end() :].strip()
        out = list(parts[1:])
        if rest:
            out.insert(0, {"text": rest, "red": bool(parts[0].get("red"))})
        return date_label, out
    return "", parts


def _clean_part_text(text: str) -> str:
    return _TATWEEL_RE.sub("", (text or "").strip())


def _prepare_day_cell(cell: dict[str, Any]) -> dict[str, Any]:
    """فصل التاريخ عن نص العمل داخل خلية اليوم."""
    prepared = dict(cell)
    parts = [
        {"text": _clean_part_text(p.get("text") or ""), "red": bool(p.get("red"))}
        for p in list(cell.get("parts") or [])
        if _clean_part_text(p.get("text") or "")
    ]
    date_label = (cell.get("date") or "").strip()
    if date_label:
        m = _DATE_RE.match(_TATWEEL_RE.sub("", date_label))
        if m:
            date_label = _format_date_label(m.group(1), m.group(2))
    if not date_label:
        date_label, parts = _split_date_and_parts(parts)
    prepared["date"] = date_label
    prepared["parts"] = parts
    prepared["is_bar"] = False
    prepared["colspan"] = 1
    return prepared


def _date_sort_key(date_label: str) -> tuple[int, int]:
    m = _DATE_RE.match((date_label or "").strip())
    if not m:
        return (99, 99)
    # التواريخ في الجدول بصيغة شهر/يوم تقريباً (9/29 ثم 10/1).
    return (int(m.group(1)), int(m.group(2)))


def _sort_week_cells_by_date(cells: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """ترتيب خلايا الأسبوع زمنياً؛ إن تعذّر يُبقى كما هو."""
    keyed = [( _date_sort_key(c.get("date") or ""), idx, c) for idx, c in enumerate(cells)]
    if any(k == (99, 99) for k, _, _ in keyed):
        return cells
    keyed.sort(key=lambda item: (item[0], item[1]))
    return [c for _, _, c in keyed]


def _append_parts_unique(target: list[dict[str, Any]], extra: list[dict[str, Any]]) -> None:
    seen = {p.get("text", "").strip() for p in target}
    for part in extra:
        text = (part.get("text") or "").strip()
        if not text or text in seen:
            continue
        target.append({"text": text, "red": bool(part.get("red"))})
        seen.add(text)


def _shape_text_parts(shape) -> list[dict[str, Any]]:
    if not getattr(shape, "has_text_frame", False):
        return []
    parts: list[dict[str, Any]] = []
    for para in shape.text_frame.paragraphs:
        runs = list(para.runs)
        if runs:
            for run in runs:
                t = (run.text or "").strip()
                if not t:
                    continue
                parts.append({"text": t, "red": _run_is_red(run)})
        else:
            t = (para.text or "").strip()
            if t:
                parts.append({"text": t, "red": False})
    if not parts:
        t = (shape.text or "").strip()
        if t:
            parts.append({"text": t, "red": False})
    return parts


def _iter_slide_text_shapes(shapes, *, skip: set[int] | None = None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    skip = skip or set()
    for shape in shapes:
        sid = id(shape)
        if sid in skip:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_slide_text_shapes(shape.shapes, skip=skip)
            continue
        if getattr(shape, "has_table", False):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        parts = _shape_text_parts(shape)
        if parts:
            yield shape, parts


def _table_cell_rects(table_shape) -> list[dict[str, Any]]:
    table = table_shape.table
    left = int(table_shape.left)
    top = int(table_shape.top)
    rects: list[dict[str, Any]] = []
    y = top
    for ri in range(len(table.rows)):
        x = left
        row_h = int(table.rows[ri].height)
        for ci in range(len(table.columns)):
            col_w = int(table.columns[ci].width)
            rects.append(
                {
                    "ri": ri,
                    "ci": ci,
                    "x1": x,
                    "y1": y,
                    "x2": x + col_w,
                    "y2": y + row_h,
                }
            )
            x += col_w
        y += row_h
    return rects


def _shape_center(shape) -> tuple[int, int]:
    return (
        int(shape.left + shape.width // 2),
        int(shape.top + shape.height // 2),
    )


def _hit_cell(rects: list[dict[str, Any]], px: int, py: int) -> dict[str, Any] | None:
    for rect in rects:
        if rect["x1"] <= px <= rect["x2"] and rect["y1"] <= py <= rect["y2"]:
            return rect
    return None


def _row_midpoints(rects: list[dict[str, Any]], nrows: int) -> list[int]:
    mids: list[int] = []
    for ri in range(nrows):
        row_rects = [r for r in rects if r["ri"] == ri]
        if not row_rects:
            mids.append(0)
            continue
        y1 = min(r["y1"] for r in row_rects)
        y2 = max(r["y2"] for r in row_rects)
        mids.append((y1 + y2) // 2)
    return mids


def _apply_overlay_shapes(
    slide,
    table_shape,
    data: dict[str, Any],
    *,
    header_rows: int,
) -> None:
    table = table_shape.table
    ncols = len(table.columns)
    rects = _table_cell_rects(table_shape)
    row_mids = _row_midpoints(rects, len(table.rows))
    rows: list[dict[str, Any]] = list(data.get("rows") or [])
    if not rows:
        return

    # خريطة (ri_body, ci) -> خلية في هيكل JSON (ri_body بدون صف العناوين)
    cell_map: dict[tuple[int, int], dict[str, Any]] = {}
    for body_ri, row in enumerate(rows):
        table_ri = body_ri + header_rows
        ci_cursor = 0
        for cell in row.get("cells") or []:
            while (table_ri, ci_cursor) in cell_map:
                ci_cursor += 1
            if ci_cursor >= ncols:
                break
            colspan = int(cell.get("colspan") or 1)
            for dc in range(colspan):
                cell_map[(table_ri, ci_cursor + dc)] = cell
            ci_cursor += colspan

    skip_ids = {id(table_shape)}
    try:
        if slide.shapes.title is not None:
            skip_ids.add(id(slide.shapes.title))
    except Exception:
        pass

    extra_bars: list[tuple[int, dict[str, Any]]] = []

    for shape, parts in _iter_slide_text_shapes(slide.shapes, skip=skip_ids):
        plain = _parts_to_plain(parts)
        if not plain:
            continue
        if _normalize_day(plain.split("\n", 1)[0]) and any(
            _normalize_day(plain.split("\n", 1)[0]) == _normalize_day(d) for d in _AR_DAYS
        ):
            continue
        if _normalize_day(plain) in {_normalize_day("برنامج التمرين"), _normalize_day("البرنامج")}:
            continue

        cx, cy = _shape_center(shape)
        wide = int(shape.width) >= int(table_shape.width * 0.55)
        hit = _hit_cell(rects, cx, cy)

        if wide:
            table_ri = 0
            for ri, mid in enumerate(row_mids):
                if cy >= mid:
                    table_ri = ri
            body_ri = table_ri - header_rows
            bar_cell = {
                "colspan": ncols,
                "rowspan": 1,
                "date": "",
                "parts": parts,
                "is_bar": True,
            }
            if body_ri < 0:
                extra_bars.append((0, bar_cell))
            elif body_ri < len(rows):
                row_cells = list(rows[body_ri].get("cells") or [])
                bar_idx = next(
                    (i for i, c in enumerate(row_cells) if c.get("is_bar")),
                    None,
                )
                if bar_idx is None:
                    row_cells.insert(0, bar_cell)
                else:
                    _append_parts_unique(row_cells[bar_idx].setdefault("parts", []), parts)
                rows[body_ri]["cells"] = row_cells
            else:
                extra_bars.append((len(rows), bar_cell))
            continue

        if not hit:
            continue

        table_ri = int(hit["ri"])
        ci = int(hit["ci"])
        body_ri = table_ri - header_rows
        if body_ri < 0:
            continue

        date_label, content_parts = _split_date_and_parts(parts)
        if _DATE_RE.match(plain):
            date_label = plain

        target = cell_map.get((table_ri, ci))
        if target is None and body_ri < len(rows):
            row_cells = rows[body_ri].get("cells") or []
            ci_pos = 0
            for cell in row_cells:
                if ci_pos <= ci < ci_pos + int(cell.get("colspan") or 1):
                    target = cell
                    break
                ci_pos += int(cell.get("colspan") or 1)

        if target is None:
            continue

        if date_label and not (target.get("date") or "").strip():
            target["date"] = date_label
        if content_parts:
            _append_parts_unique(target.setdefault("parts", []), content_parts)
        elif not date_label:
            _append_parts_unique(target.setdefault("parts", []), parts)

        if int(target.get("colspan") or 1) >= ncols and _parts_to_plain(target.get("parts") or []):
            target["is_bar"] = True

    for insert_at, bar_cell in sorted(extra_bars, key=lambda x: x[0]):
        insert_at = max(0, min(insert_at, len(rows)))
        rows.insert(insert_at, {"week_band": "", "cells": [bar_cell]})

    data["rows"] = rows


def _parts_to_plain(parts: list[dict[str, Any]]) -> str:
    return " ".join(p["text"] for p in parts if p.get("text")).strip()


def _extract_table_dict(slide, table_shape) -> dict[str, Any] | None:
    if not getattr(table_shape, "has_table", False):
        return None
    table = table_shape.table
    nrows = len(table.rows)
    ncols = len(table.columns)
    if nrows < 2 or ncols < 2:
        return None

    occupied = [[False] * ncols for _ in range(nrows)]
    header: list[str] = []
    body_rows: list[dict[str, Any]] = []

    for ri in range(nrows):
        row_cells: list[dict[str, Any]] = []
        row_band = ""
        for ci in range(ncols):
            if occupied[ri][ci]:
                continue
            cell = table.cell(ri, ci)
            if cell.is_spanned:
                continue
            rowspan = _cell_rowspan(table, ri, ci)
            colspan = _cell_colspan(cell)
            for dr in range(rowspan):
                for dc in range(colspan):
                    if ri + dr < nrows and ci + dc < ncols:
                        occupied[ri + dr][ci + dc] = True

            parts = _cell_text_parts(cell)
            date_label, content_parts = _split_date_and_parts(parts)
            plain = _parts_to_plain(content_parts if date_label else parts)
            fill_rgb = _cell_fill_rgb(cell)
            band = _classify_week_band(fill_rgb)
            if band and not row_band:
                row_band = band

            is_bar = colspan >= ncols and plain != ""
            cell_data: dict[str, Any] = {
                "colspan": colspan,
                "rowspan": rowspan,
                "date": date_label,
                "parts": content_parts if date_label else parts,
                "is_bar": is_bar,
            }
            row_cells.append(cell_data)

            if ri == 0 and plain and not header:
                norm = _normalize_day(plain)
                if any(_normalize_day(d) == norm for d in _AR_DAYS):
                    header = [(table.cell(0, c).text or "").strip() for c in range(ncols)]

        if ri == 0 and not header:
            texts = [(table.cell(0, c).text or "").strip() for c in range(ncols)]
            if sum(1 for t in texts if any(_normalize_day(t) == _normalize_day(d) for d in _AR_DAYS)) >= 4:
                header = texts

        if ri == 0 and header:
            continue
        if row_cells:
            body_rows.append({"week_band": row_band, "cells": row_cells})

    if not body_rows:
        return None
    header_rows = 1 if header else 0
    data = {"version": 1, "title": "برنامج التمرين", "header": header, "rows": body_rows}
    _apply_overlay_shapes(slide, table_shape, data, header_rows=header_rows)
    return data


def _largest_table_on_slide(slide) -> dict[str, Any] | None:
    best_shape = None
    best_score = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        tbl = shape.table
        score = len(tbl.rows) * len(tbl.columns)
        if score > best_score:
            best_shape = shape
            best_score = score
    if not best_shape:
        return None
    return _extract_table_dict(slide, best_shape)


def extract_program_table_from_slide(slide) -> dict[str, Any] | None:
    return _largest_table_on_slide(slide)


def dumps_program_table(data: dict[str, Any]) -> str:
    return json.dumps(data, ensure_ascii=False)


def loads_program_table(raw: str) -> dict[str, Any] | None:
    s = (raw or "").strip()
    if not s:
        return None
    try:
        data = json.loads(s)
    except (TypeError, ValueError):
        return None
    if not isinstance(data, dict) or not data.get("rows"):
        return None
    return data


def _render_parts_html(parts: list[dict[str, Any]]) -> str:
    chunks: list[str] = []
    for part in parts:
        text = html.escape((part.get("text") or "").strip())
        if not text:
            continue
        if part.get("red"):
            chunks.append(f'<span class="exercise-program-text-red">{text}</span>')
        else:
            chunks.append(text)
    return " ".join(chunks)


# ترتيب الأيام المطلوب: الاثنين أولاً (يمين الجدول في RTL) حتى الأحد.
_CANONICAL_DAY_ORDER = (
    "الاثنين",
    "الثلاثاء",
    "الاربعاء",
    "الخميس",
    "الجمعة",
    "السبت",
    "الاحد",
)


def _canonical_day_index(label: str) -> int | None:
    norm = _normalize_day(label)
    for idx, day in enumerate(_CANONICAL_DAY_ORDER):
        if norm == _normalize_day(day):
            return idx
    return None


def _reorder_days_monday_first(
    header: list[str], rows: list[dict[str, Any]], ncols: int
) -> tuple[list[str], list[dict[str, Any]]]:
    """إعادة ترتيب أعمدة الأيام لتبدأ بالاثنين من اليمين."""
    if len(header) != ncols:
        return header, rows
    day_indices = [_canonical_day_index(h) for h in header]
    if any(idx is None for idx in day_indices):
        return header, rows
    # ترتيب الأعمدة تصاعدياً حسب اليوم القانوني (الاثنين=0 يظهر أولاً في RTL).
    order = sorted(range(ncols), key=lambda c: day_indices[c])
    if order == list(range(ncols)):
        return header, rows

    new_header = [header[c] for c in order]
    new_rows: list[dict[str, Any]] = []
    for row in rows:
        cells = list(row.get("cells") or [])
        is_bar_row = bool(cells) and all(cell.get("is_bar") for cell in cells)
        if is_bar_row or len(cells) != ncols:
            new_rows.append(row)
            continue
        row["cells"] = [cells[c] for c in order]
        new_rows.append(row)
    return new_header, new_rows


def _prepare_bar_cell(cell: dict[str, Any], ncols: int) -> dict[str, Any]:
    prepared = dict(cell)
    parts = [
        {"text": _clean_part_text(p.get("text") or ""), "red": bool(p.get("red"))}
        for p in list(cell.get("parts") or [])
        if _clean_part_text(p.get("text") or "")
    ]
    prepared["parts"] = parts
    prepared["is_bar"] = True
    prepared["colspan"] = max(int(cell.get("colspan") or ncols), ncols)
    prepared["date"] = ""
    return prepared


_ORDER_TEXT_RE = re.compile(r"صرف\s*(?:ال)?أ?مر")


def _bar_role(cell: dict[str, Any]) -> str:
    """title = صف عنوان، week = شريط بعرض الأسبوع، phase = شريط مرحلة جزئي."""
    text = _TATWEEL_RE.sub("", _parts_to_plain(list(cell.get("parts") or [])))
    if text.startswith("ملخص") or "ملخص تمرين" in text[:24]:
        return "title"
    if text.startswith("تقييم الفصائل") or "تقييم الفصائل والسرايا" in text:
        return "week"
    if "الاستعداد لتقييم" in text or text.startswith("برنامج التمرين"):
        return "week"
    return "phase"


def _split_order_and_task(
    parts: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """فصل أوامر الصرف القصيرة عن بقية عمل اليوم."""
    order_parts: list[dict[str, Any]] = []
    task_parts: list[dict[str, Any]] = []
    for part in parts:
        text = _clean_part_text(part.get("text") or "")
        if not text:
            continue
        if _ORDER_TEXT_RE.search(text) and len(text) <= 60:
            order_parts.append({"text": text, "red": True})
            continue
        if _ORDER_TEXT_RE.search(text) and len(text) > 60:
            m = _ORDER_TEXT_RE.search(text)
            before = text[: m.start()].strip(" +|/،,")
            after = text[m.start() :].strip()
            if before:
                task_parts.append({"text": before, "red": bool(part.get("red"))})
            if after:
                order_parts.append({"text": after, "red": True})
            continue
        task_parts.append({"text": text, "red": bool(part.get("red"))})
    return order_parts, task_parts


def _slots_to_cells(slots: list[Any], ncols: int) -> list[dict[str, Any]]:
    cells: list[dict[str, Any]] = []
    i = 0
    while i < ncols:
        slot = slots[i] if i < len(slots) else None
        if slot == "SKIP":
            i += 1
            continue
        if slot is None:
            cells.append(
                {
                    "colspan": 1,
                    "rowspan": 1,
                    "date": "",
                    "parts": [],
                    "is_bar": False,
                }
            )
            i += 1
            continue
        span = max(1, int(slot.get("colspan") or 1))
        cells.append(slot)
        i += span
    return cells


def _place_span(
    slots: list[Any],
    occupied: list[bool],
    start: int,
    length: int,
    parts: list[dict[str, Any]],
) -> None:
    slots[start] = {
        "colspan": length,
        "rowspan": 1,
        "date": "",
        "parts": parts,
        "is_bar": False,
        "merged_work": True,
    }
    for i in range(start, start + length):
        occupied[i] = True
        if i > start:
            slots[i] = "SKIP"


def _bar_span(
    text: str, ncols: int
) -> tuple[int, int] | None:
    """مدى الشريط: أشرطة المراحل تمتد الإثنين → الأحد."""
    t = _TATWEEL_RE.sub("", text or "")
    if "ملخص" in t[:12]:
        return None
    # رفع الاستعداد / الاستعداد للتقييم / تقييم الفصائل / استلام المشبهات
    return (0, ncols)


def _expand_task_neighbors(detail_slots: list[Any], ncols: int) -> None:
    """دمج محدود للمهام وفق الصورة (استلام يومين، دفاع سريع+تسليم يومين)."""
    i = 0
    while i < ncols:
        slot = detail_slots[i]
        if slot in (None, "SKIP") or not isinstance(slot, dict):
            i += 1
            continue
        parts = list(slot.get("parts") or [])
        if not parts or slot.get("merged_work"):
            i += 1
            continue
        text = _parts_to_plain(parts)
        if text.startswith("اعاد") or "تنظيم والتمركز" in text:
            i += 1
            continue
        if _ORDER_TEXT_RE.search(text) and "استلام" not in text:
            i += 1
            continue

        left, right = i, i
        if text.startswith("استلام"):
            # يُعالج كشريط مرحلة منفصل؛ لا يُوسَّع هنا
            i += 1
            continue
        elif ("الدفاع السريع" in text and "تسليم" in text) and i > 0 and detail_slots[i - 1] is None:
            left = i - 1
        else:
            i += 1
            continue

        length = right - left + 1
        # للإاستلام: النص بدون أمر الصرف (الأمر يبقى في يومه إن أمكن)
        use_parts = parts
        if text.startswith("استلام"):
            _o, task_parts = _split_order_and_task(parts)
            use_parts = task_parts or parts
            # أعد أمر الصرف على يوم الاستلام الأصلي إن دُمج لليمين
            if _o and left == i and right > i:
                # الأمر على نفس بداية الدمج (الثلاثاء)
                pass
            elif _o:
                # احتفظ بالأمر مع المهمة في الخلية المدمجة
                use_parts = task_parts + _o

        for j in range(left, right + 1):
            detail_slots[j] = None
        detail_slots[left] = {
            "colspan": length,
            "rowspan": 1,
            "date": "",
            "parts": use_parts,
            "is_bar": False,
            "merged_work": True,
        }
        for j in range(left + 1, right + 1):
            detail_slots[j] = "SKIP"
        # إن وُجد أمر منفصل بعد تقسيم استلام، ضعه على الثلاثاء داخل النص المدمج
        i = right + 1


def _build_week_layers(
    week_cells: list[dict[str, Any]],
    bars: list[dict[str, Any]],
    ncols: int,
) -> tuple[list[list[dict[str, Any]]], list[dict[str, Any]], list[dict[str, Any]]]:
    """طبقات الأسبوع مطابقة للصورة: أشرطة مرحلة (قد تتعدد) + صناديق يومية."""
    detail_slots: list[Any] = [None] * ncols
    istelam_parts: list[dict[str, Any]] | None = None

    for idx, cell in enumerate(week_cells[:ncols]):
        order_parts, task_parts = _split_order_and_task(list(cell.get("parts") or []))
        if task_parts and _parts_to_plain(task_parts).startswith("استلام"):
            istelam_parts = task_parts
            if order_parts:
                detail_slots[idx] = {
                    "colspan": 1,
                    "rowspan": 1,
                    "date": "",
                    "parts": order_parts,
                    "is_bar": False,
                }
            continue
        if task_parts or order_parts:
            detail_slots[idx] = {
                "colspan": 1,
                "rowspan": 1,
                "date": "",
                "parts": task_parts + order_parts,
                "is_bar": False,
            }

    _expand_task_neighbors(detail_slots, ncols)

    # صفوف مرحلة متعددة (شريط طويل + شريط قصير فوقه كما في الصورة)
    phase_layers: list[list[Any]] = []
    phase_occ: list[list[bool]] = []

    def _add_phase_bar(start: int, length: int, parts: list[dict[str, Any]]) -> None:
        for slots, occupied in zip(phase_layers, phase_occ):
            if not any(occupied[i] for i in range(start, start + length)):
                _place_span(slots, occupied, start, length, parts)
                return
        slots = [None] * ncols
        occupied = [False] * ncols
        _place_span(slots, occupied, start, length, parts)
        phase_layers.append(slots)
        phase_occ.append(occupied)

    scored: list[tuple[int, list[dict[str, Any]], str]] = []
    for bar in bars:
        text = _parts_to_plain(list(bar.get("parts") or []))
        span = _bar_span(text, ncols)
        if span is None:
            continue
        scored.append((span[1], list(bar.get("parts") or []), text))
    if istelam_parts:
        span = _bar_span("استلام", ncols)
        if span:
            scored.append((span[1], istelam_parts, "استلام"))
    # الأطول أولاً
    scored.sort(key=lambda item: -item[0])
    for _ln, parts, _text in scored:
        span = _bar_span(_parts_to_plain(parts) if not _text.startswith("استلام") else "استلام", ncols)
        if not span:
            continue
        s, ln = span
        _add_phase_bar(s, ln, parts)

    phase_rows = [_slots_to_cells(slots, ncols) for slots in phase_layers]
    return phase_rows, _slots_to_cells(detail_slots, ncols), []


def _normalize_calendar_rows(
    rows: list[dict[str, Any]], ncols: int
) -> list[dict[str, Any]]:
    """بناء أسابيع جدول الأعمال مطابقة للصورة المرجعية."""
    weeks: list[dict[str, Any]] = []
    leading_bars: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None

    for row in rows:
        cells = list(row.get("cells") or [])
        bars = [_prepare_bar_cell(cell, ncols) for cell in cells if cell.get("is_bar")]
        date_cells = [
            _prepare_day_cell(cell) for cell in cells if not cell.get("is_bar")
        ]

        if date_cells:
            week_cells = _sort_week_cells_by_date(date_cells)[:ncols]
            while len(week_cells) < ncols:
                week_cells.append(
                    {
                        "colspan": 1,
                        "rowspan": 1,
                        "date": "",
                        "parts": [],
                        "is_bar": False,
                    }
                )
            current = {
                "cells": week_cells,
                "bars": list(leading_bars) + bars,
            }
            leading_bars = []
            weeks.append(current)
            continue

        if current is not None:
            current["bars"].extend(bars)
        else:
            leading_bars.extend(bars)

    if leading_bars and weeks:
        weeks[0]["bars"] = leading_bars + list(weeks[0].get("bars") or [])

    normalized: list[dict[str, Any]] = []
    week_bands = ("beige", "blue", "blue", "beige", "")
    for week_index, week in enumerate(weeks):
        band = week_bands[week_index] if week_index < len(week_bands) else ""
        week_cells = list(week.get("cells") or [])
        bars = list(week.get("bars") or [])
        phase_rows, detail_cells, titles = _build_week_layers(week_cells, bars, ncols)

        normalized.append(
            {
                "week_band": band,
                "row_kind": "dates",
                "cells": [
                    {
                        "colspan": 1,
                        "rowspan": 1,
                        "date": (c.get("date") or "").strip(),
                        "parts": [],
                        "is_bar": False,
                    }
                    for c in week_cells
                ],
            }
        )
        for phase_cells in phase_rows:
            if any(c.get("parts") for c in phase_cells):
                normalized.append(
                    {"week_band": band, "row_kind": "phase", "cells": phase_cells}
                )
        if any(c.get("parts") for c in detail_cells):
            normalized.append(
                {
                    "week_band": band,
                    "row_kind": "activities",
                    "cells": detail_cells,
                }
            )
        for bar in titles:
            normalized.append(
                {"week_band": band, "row_kind": "bar", "cells": [bar]}
            )

    return normalized


def _canonical_header(ncols: int) -> list[str]:
    labels = ["الإثنين", "الثلاثاء", "الأربعاء", "الخميس", "الجمعة", "السبت", "الأحد"]
    if ncols <= len(labels):
        return labels[:ncols]
    return labels + [f"يوم {i}" for i in range(len(labels) + 1, ncols + 1)]


def render_program_table_html(raw_json: str) -> str:
    data = loads_program_table(raw_json)
    if not data:
        return ""

    title = html.escape((data.get("title") or "برنامج التمرين").strip())
    header: list[str] = list(data.get("header") or [])
    rows: list[dict[str, Any]] = list(data.get("rows") or [])
    ncols = len(header) or 7
    rows = _normalize_calendar_rows(rows, ncols)
    header = _canonical_header(ncols)

    out: list[str] = [
        '<div class="exercise-program-calendar-wrap">',
        '<div class="exercise-program-calendar-heading">',
        f'<h3 class="exercise-program-calendar-title">{title}</h3>',
        "</div>",
        '<table class="exercise-program-calendar" dir="rtl">',
    ]

    out.append("<thead><tr>")
    for h in header:
        out.append(f"<th>{html.escape((h or '').strip())}</th>")
    out.append("</tr></thead>")

    out.append("<tbody>")
    for row in rows:
        band = (row.get("week_band") or "").strip()
        cells = list(row.get("cells") or [])
        row_kind = (row.get("row_kind") or "").strip()
        is_bar_row = row_kind == "bar" or (
            bool(cells) and all(cell.get("is_bar") for cell in cells)
        )
        is_dates_row = row_kind == "dates"
        is_acts_row = row_kind == "activities"
        is_phase_row = row_kind == "phase"
        tr_classes = []
        if band:
            tr_classes.append(f"week-{band}")
        if is_bar_row:
            tr_classes.append("week-bar-row")
        elif is_dates_row:
            tr_classes.append("week-dates-row")
        elif is_phase_row:
            tr_classes.append("week-phase-row")
        elif is_acts_row:
            tr_classes.append("week-acts-row")
        else:
            tr_classes.append("week-days-row")
        tr_cls = f' class="{" ".join(tr_classes)}"' if tr_classes else ""
        out.append(f"<tr{tr_cls}>")
        for cell in cells:
            colspan = int(cell.get("colspan") or 1)
            rowspan = int(cell.get("rowspan") or 1)
            attrs = []
            if colspan > 1:
                attrs.append(f'colspan="{colspan}"')
            if rowspan > 1:
                attrs.append(f'rowspan="{rowspan}"')
            is_bar = bool(cell.get("is_bar"))
            merged_work = bool(cell.get("merged_work"))
            classes = []
            if is_bar:
                classes.append("exercise-program-bar-cell")
            elif is_dates_row:
                classes.append("exercise-program-date-cell")
            else:
                classes.append("exercise-program-day-cell")
            if merged_work:
                classes.append("is-merged-work")
            if classes:
                attrs.append(f'class="{" ".join(classes)}"')
            attr_s = (" " + " ".join(attrs)) if attrs else ""
            date_label = html.escape((cell.get("date") or "").strip())
            inner = _render_parts_html(list(cell.get("parts") or []))
            if is_bar:
                out.append(
                    f'<td{attr_s}><div class="exercise-program-bar-box">{inner}</div></td>'
                )
            elif is_dates_row:
                out.append(
                    f'<td{attr_s}><div class="exercise-program-cell-date">{date_label}</div></td>'
                )
            else:
                box = (
                    f'<div class="exercise-program-cell-box">{inner}</div>'
                    if inner
                    else ""
                )
                out.append(f"<td{attr_s}>{box}</td>")
        out.append("</tr>")
    out.append("</tbody></table></div>")
    return "".join(out)
