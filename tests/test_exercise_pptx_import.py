"""اختبار استيراد تبويبات التمرين من PowerPoint."""

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.exercise_pptx_import import parse_exercise_pptx_bytes
from app.exercise_program_table import (
    extract_program_table_from_slide,
    loads_program_table,
    render_program_table_html,
)


def _build_sample_pptx() -> bytes:
    prs = Presentation()
    slides = [
        ("القصد", "قصد التمرين المستورد من الملف."),
        ("المشاركون في التمرين", "لواء 1\nلواء 2"),
        ("الفكرة العامة", "نص الفكرة العامة للتمرين."),
        ("الفكرة الخاصة", "نص الفكرة الخاصة."),
        ("البرنامج", "جدول البرنامج التفصيلي."),
        ("الخريطة", "وصف الخريطة."),
    ]
    for title, body in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_info_slide_pptx() -> bytes:
    """شريحة «معلومات التمرين» بعناوين فرعية للقصد والمشاركين."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "معلومات التمرين"
    slide.placeholders[1].text = (
        "القصد\n"
        "تحقيق الجاهزية القتالية.\n\n"
        "المشاركون في التمرين\n"
        "مجموعة لواء 1\n"
        "وحدات الدعم"
    )
    for title, body in [
        ("الفكرة العامة", "نص عام."),
        ("الفكرة الخاصة", "نص خاص."),
        ("البرنامج", "برنامج."),
        ("الخريطة", "خريطة."),
    ]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        s.placeholders[1].text = body
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_program_table_pptx() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_box.text_frame.text = "البرنامج"

    days = ["الإثنين", "الثلاثاء", "الأربعاء", "الخميس", "الجمعة", "السبت", "الأحد"]
    rows, cols = 3, 7
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.25), Inches(0.85), Inches(9.5), Inches(2.2)
    )
    table = table_shape.table
    for ci, day in enumerate(days):
        cell = table.cell(0, ci)
        cell.text = day
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(228, 228, 228)

    bar = table.cell(1, 0)
    bar.merge(table.cell(1, 6))
    bar.text = "رفع حالة الاستعداد وتفتيش الجاهزية"
    for ci in range(cols):
        c = table.cell(1, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(235, 232, 228)

    warn = table.cell(2, 0)
    warn.text = "9/29"
    p = warn.text_frame.paragraphs[0]
    p.clear()
    run_date = p.add_run()
    run_date.text = "9/29"
    run_red = p.add_run()
    run_red.text = " صرف الأمر الإنذاري/1"
    run_red.font.color.rgb = RGBColor(192, 0, 0)
    for ci in range(cols):
        c = table.cell(2, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(245, 240, 230)

    # شرائح أخرى
    for title, body in [
        ("الفكرة العامة", "نص عام."),
        ("الفكرة الخاصة", "نص خاص."),
        ("الخريطة", "خريطة."),
    ]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        s.placeholders[1].text = body

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_program_overlay_pptx() -> bytes:
    """جدول بتواريخ فقط في الخلايا والمحتوى في أشكال منفصلة (كما في PowerPoint الحقيقي)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_box.text_frame.text = "البرنامج"

    days = ["الإثنين", "الثلاثاء", "الأربعاء", "الخميس", "الجمعة", "السبت", "الأحد"]
    rows, cols = 2, 7
    table_top = Inches(0.85)
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.25), table_top, Inches(9.5), Inches(1.4)
    )
    table = table_shape.table
    for ci, day in enumerate(days):
        table.cell(0, ci).text = day
    dates = ["9/29", "9/30", "10/1", "10/2", "10/3", "10/4", "10/5"]
    row_h = int(table.rows[1].height)
    col_w = int(table.columns[0].width)
    tbl_left = int(table_shape.left)
    tbl_top = int(table_shape.top)
    header_h = int(table.rows[0].height)
    for ci, date in enumerate(dates):
        table.cell(1, ci).text = date
        box = slide.shapes.add_textbox(
            tbl_left + ci * col_w + int(col_w * 0.05),
            tbl_top + header_h + int(row_h * 0.35),
            int(col_w * 0.9),
            int(row_h * 0.55),
        )
        tf = box.text_frame
        tf.clear()
        if ci == 0:
            run = tf.paragraphs[0].add_run()
            run.text = "صرف الأمر الإنذاري/1"
            run.font.color.rgb = RGBColor(192, 0, 0)
        elif ci == 6:
            tf.paragraphs[0].text = "إعادة التنظيم والتمركز في منطقة التدريب"

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_parse_program_table_overlay_shapes():
    data = _build_program_overlay_pptx()
    out = parse_exercise_pptx_bytes(data)
    assert out["ok"] is True
    table = loads_program_table(out["fields"]["program_table_json"])
    assert table is not None
    first_row = table["rows"][0]
    first_cell = first_row["cells"][0]
    assert first_cell.get("date") == "9/29"
    parts_text = " ".join(p["text"] for p in first_cell.get("parts") or [])
    assert "صرف الأمر الإنذاري" in parts_text
    last_cell = first_row["cells"][-1]
    last_text = " ".join(p["text"] for p in last_cell.get("parts") or [])
    assert "التمركز" in last_text or "التنظيم" in last_text


def test_parse_exercise_pptx_by_slide_titles():
    data = _build_sample_pptx()
    out = parse_exercise_pptx_bytes(data)
    assert out["ok"] is True
    fields = out["fields"]
    assert "قصد التمرين" in fields["exercise_purpose"]
    assert "لواء 1" in fields["exercise_participants"]
    assert "الفكرة العامة" in fields["general_idea_text"] or "نص الفكرة العامة" in fields["general_idea_text"]
    assert "الفكرة الخاصة" in fields["specific_idea_text"] or "نص الفكرة الخاصة" in fields["specific_idea_text"]
    assert "البرنامج" in fields["program_text"] or "جدول البرنامج" in fields["program_text"]
    assert "الخريطة" in fields["map_text"] or "وصف الخريطة" in fields["map_text"]


def test_parse_info_fields_from_info_slide():
    data = _build_info_slide_pptx()
    out = parse_exercise_pptx_bytes(data)
    assert out["ok"] is True
    fields = out["fields"]
    assert "الجاهزية" in fields["exercise_purpose"]
    assert "مجموعة لواء 1" in fields["exercise_participants"]
    assert "وحدات الدعم" in fields["exercise_participants"]


def test_parse_program_table_from_pptx():
    data = _build_program_table_pptx()
    out = parse_exercise_pptx_bytes(data)
    assert out["ok"] is True
    fields = out["fields"]
    assert fields.get("program_table_json")
    table = loads_program_table(fields["program_table_json"])
    assert table is not None
    assert len(table.get("header") or []) == 7
    assert len(table.get("rows") or []) >= 2
    html_out = out.get("program_table_html") or render_program_table_html(fields["program_table_json"])
    assert "exercise-program-calendar" in html_out
    assert "صرف الأمر الإنذاري" in html_out
    assert "exercise-program-text-red" in html_out


if __name__ == "__main__":
    test_parse_exercise_pptx_by_slide_titles()
    test_parse_info_fields_from_info_slide()
    test_parse_program_table_from_pptx()
    test_parse_program_table_overlay_shapes()
    print("ok")
